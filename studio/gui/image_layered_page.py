# -*- coding: utf-8 -*-
import os
import sys
import shutil
import zipfile
import tempfile
import random
import traceback
from PIL import Image

from PySide6.QtWidgets import (QVBoxLayout, QHBoxLayout, QLabel, QPushButton,
                               QFileDialog, QProgressBar, QMessageBox, QFrame,
                               QSlider, QSplitter, QWidget, QStackedWidget,
                               QButtonGroup, QLineEdit, QScrollArea, QGridLayout, QDialog,
                               QCheckBox)
from PySide6.QtCore import Signal, QThread, Qt, QRectF, QPoint
from utils.base_worker import BaseWorker
from PySide6.QtGui import QImage, QPixmap, QPainter, QColor
from utils.logger_utils import log
from config.paths import TMP_DIR, QWEN_IMAGE_LAYERED_DIR

_cached_pipeline = None

class QwenLayeredWorker(BaseWorker):
    progress = Signal(str)
    finished = Signal(bool, list, str, str, str, str) # success, layer_paths, pptx_path, zip_path, psd_path, error_msg
    
    def __init__(self, img_path, layers_count=4, num_steps=50, scale=4.0, prompt="", neg_prompt=" ", seed=-1, low_vram=True):
        super().__init__()
        self.img_path = img_path
        self.layers_count = layers_count
        self.num_steps = num_steps
        self.scale = scale
        self.prompt = prompt if prompt else None
        self.neg_prompt = neg_prompt
        self.seed = seed
        self.low_vram = low_vram
        
    def run(self):
        try:
            self.progress.emit("检查 GPU 加速硬件支持...")
            import torch
            torch.cuda.empty_cache()
            if not torch.cuda.is_available():
                raise RuntimeError("当前系统缺少 NVIDIA CUDA 支持。图像智能分层模型体量巨大，必须在 NVIDIA GPU 下执行！")
                
            self.progress.emit("正在构建 Qwen-Image-Layered 运行环境...")
            import sys
            import os
            
            # Add Qwen repository path dynamically
            qwen_repo_dir = QWEN_IMAGE_LAYERED_DIR
            if qwen_repo_dir not in sys.path:
                sys.path.insert(0, qwen_repo_dir)
                
            from diffusers import QwenImageLayeredPipeline
            from pptx import Presentation
            from psd_tools import PSDImage
            
            global _cached_pipeline
            if _cached_pipeline is None:
                if os.path.exists(os.path.join(qwen_repo_dir, "model_index.json")):
                    self.progress.emit("正在从本地目录 (apps/Qwen-Image-Layered) 载入 Qwen-Image-Layered 预训练模型...")
                    model_path = qwen_repo_dir
                elif os.path.exists(os.path.join(os.path.dirname(qwen_repo_dir), "Qwen-Image-Layered-Weights", "model_index.json")):
                    self.progress.emit("正在从本地目录 (apps/Qwen-Image-Layered-Weights) 载入 Qwen-Image-Layered 预训练模型...")
                    model_path = os.path.join(os.path.dirname(qwen_repo_dir), "Qwen-Image-Layered-Weights")
                else:
                    os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"
                    self.progress.emit("正在从云端下载并载入 Qwen-Image-Layered 预训练模型 (模型约 45GB，初次使用国内镜像高速下载，请保持联网)...")
                    model_path = "Qwen/Qwen-Image-Layered"
                
                _cached_pipeline = QwenImageLayeredPipeline.from_pretrained(model_path, torch_dtype=torch.bfloat16)
                _cached_pipeline.set_progress_bar_config(disable=True)

            if self.low_vram:
                self.progress.emit("启用显存优化模式 (CPU Offload)...")
                _cached_pipeline.enable_model_cpu_offload()
            else:
                self.progress.emit("载入模型到 GPU 显存...")
                _cached_pipeline = _cached_pipeline.to("cuda")
                
            # Helper: PPTX Converter
            def imagelist_to_pptx(img_files):
                with Image.open(img_files[0]) as img:
                    w_px, h_px = img.size
                def px_to_emu(px, dpi=96):
                    inch = px / dpi
                    emu = inch * 914400
                    return int(emu)
                prs = Presentation()
                prs.slide_width = px_to_emu(w_px)
                prs.slide_height = px_to_emu(h_px)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                left = top = 0
                for path in img_files:
                    slide.shapes.add_picture(path, left, top, width=px_to_emu(w_px), height=px_to_emu(h_px))
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                    prs.save(tmp.name)
                    return tmp.name
                    
            # Helper: PSD Converter
            def imagelist_to_psd(img_files):
                layers = []
                for path in img_files:
                    layers.append(Image.open(path).convert('RGBA'))
                w, h = layers[0].size
                psd = PSDImage.new(mode='RGBA', size=(w, h))
                for idx, img in enumerate(layers):
                    layer = psd.create_pixel_layer(image=img, name=f"Layer {idx + 1}")
                    psd.append(layer)
                with tempfile.NamedTemporaryFile(suffix=".psd", delete=False) as tmp:
                    psd.save(tmp.name)
                    return tmp.name
            
            # Setup seed
            actual_seed = self.seed if self.seed >= 0 else random.randint(0, 2147483647)
            
            # Format inputs
            pil_image = Image.open(self.img_path).convert("RGB").convert("RGBA")
            inputs = {
                "image": pil_image,
                "generator": torch.Generator(device='cuda').manual_seed(actual_seed),
                "true_cfg_scale": self.scale,
                "prompt": self.prompt,
                "negative_prompt": self.neg_prompt,
                "num_inference_steps": self.num_steps,
                "num_images_per_prompt": 1,
                "layers": self.layers_count,
                "resolution": 640,
                "cfg_normalize": True,
                "use_en_prompt": True,
            }
            
            self.progress.emit(f"AI 正在开始进行图层分层解耦 (图层数:{self.layers_count}, 步数:{self.num_steps})...")
            with torch.inference_mode():
                output = _cached_pipeline(**inputs)
                output_images = output.images[0]
            
            torch.cuda.empty_cache()
            self.progress.emit("分层运算完成，正在导出临时文件...")
            temp_files = []
            
            # Create permanent sub-directory in runtime for output display
            tmp_dir = os.path.join(TMP_DIR, "layered_output")
            if os.path.exists(tmp_dir):
                shutil.rmtree(tmp_dir)
            os.makedirs(tmp_dir, exist_ok=True)
            
            for idx, img in enumerate(output_images):
                layer_path = os.path.join(tmp_dir, f"layer_{idx+1}.png")
                img.save(layer_path, format="PNG")
                temp_files.append(layer_path)
                
            self.progress.emit("正在生成多图层 PPTX 文件...")
            pptx_path = imagelist_to_pptx(temp_files)
            
            self.progress.emit("正在打包图层 ZIP 压缩包...")
            zip_out_path = os.path.join(tmp_dir, "layers_archive.zip")
            with zipfile.ZipFile(zip_out_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for idx, path in enumerate(temp_files):
                    zipf.write(path, f"layer_{idx+1}.png")
                    
            self.progress.emit("正在合成 Photoshop PSD 分层文件...")
            psd_path = imagelist_to_psd(temp_files)
            
            self.finished.emit(True, temp_files, pptx_path, zip_out_path, psd_path, "")
            
        except Exception as e:
            tb = traceback.format_exc()
            log.error(f"Qwen layered process error: {e}\n{tb}")
            self.finished.emit(False, [], "", "", "", str(e))


class ResultPreviewWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.image = QImage()
        
    def set_image(self, image):
        self.image = image
        self.update()
        
    def paintEvent(self, event):
        painter = QPainter(self)
        W, H = self.width(), self.height()
        
        # Draw Dark Grid Checkerboard
        tile_size = 10
        color1 = QColor(32, 32, 34)
        color2 = QColor(44, 44, 46)
        for x in range(0, W, tile_size):
            for y in range(0, H, tile_size):
                if ((x // tile_size) + (y // tile_size)) % 2 == 0:
                    painter.fillRect(x, y, tile_size, tile_size, color1)
                else:
                    painter.fillRect(x, y, tile_size, tile_size, color2)
                    
        # Draw Layer image
        if not self.image.isNull():
            img_w, img_h = self.image.width(), self.image.height()
            factor_w = W / img_w
            factor_h = H / img_h
            scale_factor = min(factor_w, factor_h)
            
            display_w = img_w * scale_factor
            display_h = img_h * scale_factor
            offset_x = (W - display_w) / 2
            offset_y = (H - display_h) / 2
            
            display_rect = QRectF(offset_x, offset_y, display_w, display_h)
            painter.drawImage(display_rect, self.image)


class ImagePreviewWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.image = QImage()
        
    def set_image(self, image_or_path):
        if isinstance(image_or_path, str):
            self.image = QImage(image_or_path)
        else:
            self.image = image_or_path
        self.update()
        
    def paintEvent(self, event):
        painter = QPainter(self)
        painter.fillRect(self.rect(), QColor(26, 26, 28))
        
        if not self.image.isNull():
            W, H = self.width(), self.height()
            img_w, img_h = self.image.width(), self.image.height()
            
            factor_w = W / img_w
            factor_h = H / img_h
            scale_factor = min(factor_w, factor_h)
            
            display_w = img_w * scale_factor
            display_h = img_h * scale_factor
            offset_x = (W - display_w) / 2
            offset_y = (H - display_h) / 2
            
            display_rect = QRectF(offset_x, offset_y, display_w, display_h)
            painter.drawImage(display_rect, self.image)


class LayerCardWidget(QFrame):
    def __init__(self, layer_idx, img_path, on_preview, on_save):
        super().__init__()
        self.setObjectName("card")
        self.setStyleSheet("QFrame#card { background-color: #2c2c2e; border: 1px solid #3a3a3c; border-radius: 10px; }")
        
        layout = QVBoxLayout(self)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(10)
        
        # Checkerboard preview
        self.preview = ResultPreviewWidget()
        self.preview.setFixedSize(160, 160)
        self.preview.set_image(QImage(img_path))
        layout.addWidget(self.preview, 0, Qt.AlignCenter)
        
        # Layer label
        lbl_info = QLabel(f"图层 {layer_idx} (Layer)")
        lbl_info.setAlignment(Qt.AlignCenter)
        lbl_info.setStyleSheet("font-weight: bold; font-size: 13px; color: #ffffff;")
        layout.addWidget(lbl_info)
        
        # Buttons layout
        btn_layout = QHBoxLayout()
        btn_view = QPushButton("🔍 查看")
        btn_view.setObjectName("pill_button")
        btn_view.setCursor(Qt.PointingHandCursor)
        btn_view.clicked.connect(lambda: on_preview(img_path, layer_idx))
        
        btn_export = QPushButton("💾 另存")
        btn_export.setObjectName("pill_button")
        btn_export.setCursor(Qt.PointingHandCursor)
        btn_export.clicked.connect(lambda: on_save(img_path, layer_idx))
        
        btn_layout.addWidget(btn_view)
        btn_layout.addWidget(btn_export)
        layout.addLayout(btn_layout)


class LayerPreviewDialog(QDialog):
    def __init__(self, img_path, layer_idx, parent=None):
        super().__init__(parent)
        self.setWindowTitle(f"查看图层 {layer_idx} (大图)")
        self.setMinimumSize(600, 600)
        self.resize(750, 750)
        self.setStyleSheet("QDialog { background-color: #1a1a1c; }")
        
        layout = QVBoxLayout(self)
        layout.setContentsMargins(15, 15, 15, 15)
        layout.setSpacing(10)
        
        # Transparent Preview
        self.view = ResultPreviewWidget()
        self.view.set_image(QImage(img_path))
        layout.addWidget(self.view, 1)
        
        # Footer
        footer = QHBoxLayout()
        footer.addStretch()
        btn_close = QPushButton("关闭")
        btn_close.clicked.connect(self.accept)
        footer.addWidget(btn_close)
        layout.addLayout(footer)


from gui.base_page import BasePage


class ImageLayeredPage(BasePage):
    def __init__(self, parent_widget, main_window):
        super().__init__(parent_widget, main_window)
        self.parent = parent_widget  # 兼容旧引用
        self.selected_img_path = None
        self.worker = None
        
        # Paths to generated formats
        self.layer_paths = []
        self.pptx_path = None
        self.zip_path = None
        self.psd_path = None
        
    def setup(self):
        # Page layout
        main_layout = QVBoxLayout(self.parent)
        main_layout.setContentsMargins(40, 40, 40, 40)
        main_layout.setSpacing(20)
        
        # 1. Header
        header_layout = QHBoxLayout()
        heading = QLabel("AI 图像智能分层 (Qwen Image Layered)")
        heading.setObjectName("heading")
        header_layout.addWidget(heading)
        header_layout.addStretch()
        main_layout.addLayout(header_layout)

        # Check if the capability module exists
        if not os.path.exists(QWEN_IMAGE_LAYERED_DIR):
            card = QFrame()
            card.setObjectName("card")
            card_layout = QVBoxLayout(card)
            card_layout.setContentsMargins(40, 40, 40, 40)
            card_layout.setSpacing(20)
            card_layout.setAlignment(Qt.AlignCenter)
            
            icon_label = QLabel("⚠️")
            icon_label.setAlignment(Qt.AlignCenter)
            icon_label.setStyleSheet("font-size: 64px;")
            card_layout.addWidget(icon_label)
            
            msg_label = QLabel("此功能需要配置对应的能力模块")
            msg_label.setAlignment(Qt.AlignCenter)
            msg_label.setStyleSheet("font-size: 18px; font-weight: bold; color: #ff9900;")
            card_layout.addWidget(msg_label)
            
            desc_label = QLabel(
                "未在 apps 目录下检测到 Qwen-Image-Layered 模块。\n"
                "请将对应的能力模块文件夹放置于以下路径：\n"
                f"{QWEN_IMAGE_LAYERED_DIR}"
            )
            desc_label.setAlignment(Qt.AlignCenter)
            desc_label.setStyleSheet("font-size: 13px; color: #8e8e93; line-height: 20px;")
            card_layout.addWidget(desc_label)
            
            # Put the card in the center of main_layout
            center_layout = QHBoxLayout()
            center_layout.addStretch()
            center_layout.addWidget(card)
            center_layout.addStretch()
            
            main_layout.addStretch()
            main_layout.addLayout(center_layout)
            main_layout.addStretch()
            return
        
        # 2. Workspace splitter
        splitter = QSplitter(Qt.Horizontal)
        splitter.setStyleSheet("QSplitter::handle { background-color: #2e2e32; width: 2px; }")
        
        # Left Panel (Controls)
        left_panel = QFrame()
        left_panel.setObjectName("card")
        left_panel.setFixedWidth(320)
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(20, 20, 20, 20)
        left_layout.setSpacing(15)
        
        # Image Upload
        left_layout.addWidget(QLabel("1. 选择原始图像:"))
        self.btn_select_image = QPushButton("📁 上传本地图像")
        self.btn_select_image.setObjectName("secondary_button")
        self.btn_select_image.setCursor(Qt.PointingHandCursor)
        self.btn_select_image.clicked.connect(self.choose_image)
        left_layout.addWidget(self.btn_select_image)
        
        self.lbl_image_info = QLabel("未载入任何图像，请先上传。")
        self.lbl_image_info.setObjectName("muted_text")
        self.lbl_image_info.setWordWrap(True)
        left_layout.addWidget(self.lbl_image_info)
        
        left_layout.addWidget(self.create_separator())
        
        # Settings
        left_layout.addWidget(QLabel("2. 图像分解参数:"))
        
        # Layers slider
        layers_row = QHBoxLayout()
        layers_row.addWidget(QLabel("分解图层数:"))
        self.lbl_layers_val = QLabel("4")
        self.lbl_layers_val.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
        self.lbl_layers_val.setObjectName("accent_text")
        layers_row.addWidget(self.lbl_layers_val)
        left_layout.addLayout(layers_row)
        
        self.slider_layers = QSlider(Qt.Horizontal)
        self.slider_layers.setRange(2, 10)
        self.slider_layers.setValue(4)
        self.slider_layers.valueChanged.connect(self.change_layers_val)
        left_layout.addWidget(self.slider_layers)
        
        # Steps slider
        steps_row = QHBoxLayout()
        steps_row.addWidget(QLabel("推理步数 (Steps):"))
        self.lbl_steps_val = QLabel("50")
        self.lbl_steps_val.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
        self.lbl_steps_val.setObjectName("accent_text")
        steps_row.addWidget(self.lbl_steps_val)
        left_layout.addLayout(steps_row)
        
        self.slider_steps = QSlider(Qt.Horizontal)
        self.slider_steps.setRange(10, 50)
        self.slider_steps.setValue(50)
        self.slider_steps.valueChanged.connect(self.change_steps_val)
        left_layout.addWidget(self.slider_steps)
        
        # CFG scale
        cfg_row = QHBoxLayout()
        cfg_row.addWidget(QLabel("引导系数 (CFG):"))
        self.lbl_cfg_val = QLabel("4.0")
        self.lbl_cfg_val.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
        self.lbl_cfg_val.setObjectName("accent_text")
        cfg_row.addWidget(self.lbl_cfg_val)
        left_layout.addLayout(cfg_row)
        
        self.slider_cfg = QSlider(Qt.Horizontal)
        self.slider_cfg.setRange(10, 100) # Mapping to 1.0 - 10.0
        self.slider_cfg.setValue(40)
        self.slider_cfg.valueChanged.connect(self.change_cfg_val)
        left_layout.addWidget(self.slider_cfg)
        
        # Prompt (Optional)
        left_layout.addWidget(QLabel("图像描述 (Prompt - 可选):"))
        self.txt_prompt = QLineEdit()
        self.txt_prompt.setPlaceholderText("输入描述以辅助识别遮挡关系 (例如: 一个女孩站在海滩上)")
        left_layout.addWidget(self.txt_prompt)
        
        # Low VRAM Optimization Checkbox
        self.chk_low_vram = QCheckBox("启用显存优化 (低显存/防 OOM 模式)")
        self.chk_low_vram.setChecked(True)
        self.chk_low_vram.setToolTip("开启后将启用 CPU Offload 机制，仅在计算时将必要图层载入显存，大幅降低显存需求，防止 24G/16G/12G/8G 显卡发生 Out of Memory 错误")
        left_layout.addWidget(self.chk_low_vram)
        
        left_layout.addWidget(self.create_separator())
        
        # Execution & Actions
        left_layout.addWidget(QLabel("3. 分层操作:"))
        self.lbl_status = QLabel("")
        self.lbl_status.setObjectName("muted_text")
        left_layout.addWidget(self.lbl_status)
        
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        left_layout.addWidget(self.progress_bar)
        
        self.btn_run = QPushButton("⚡ 开始智能分层")
        self.btn_run.setObjectName("primary_button")
        self.btn_run.setCursor(Qt.PointingHandCursor)
        self.btn_run.setEnabled(False)
        self.btn_run.clicked.connect(self.run_layering)
        left_layout.addWidget(self.btn_run)
        
        # Export Actions Group
        left_layout.addWidget(QLabel("4. 导出整图工程:"))
        
        export_layout = QHBoxLayout()
        self.btn_export_pptx = QPushButton("📊 PPTX")
        self.btn_export_pptx.setObjectName("secondary_button")
        self.btn_export_pptx.setCursor(Qt.PointingHandCursor)
        self.btn_export_pptx.setEnabled(False)
        self.btn_export_pptx.clicked.connect(self.save_pptx)
        
        self.btn_export_zip = QPushButton("📦 ZIP")
        self.btn_export_zip.setObjectName("secondary_button")
        self.btn_export_zip.setCursor(Qt.PointingHandCursor)
        self.btn_export_zip.setEnabled(False)
        self.btn_export_zip.clicked.connect(self.save_zip)
        
        self.btn_export_psd = QPushButton("🎨 PSD")
        self.btn_export_psd.setObjectName("secondary_button")
        self.btn_export_psd.setCursor(Qt.PointingHandCursor)
        self.btn_export_psd.setEnabled(False)
        self.btn_export_psd.clicked.connect(self.save_psd)
        
        export_layout.addWidget(self.btn_export_pptx)
        export_layout.addWidget(self.btn_export_zip)
        export_layout.addWidget(self.btn_export_psd)
        left_layout.addLayout(export_layout)
        
        left_layout.addStretch()
        splitter.addWidget(left_panel)
        
        # Right Panel (Gallery & Previews)
        right_panel = QFrame()
        right_panel.setObjectName("card")
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(15, 15, 15, 15)
        right_layout.setSpacing(12)
        
        # Preview view headers
        view_tabs_layout = QHBoxLayout()
        self.btn_view_orig = QPushButton("🎨 原始图像 (Original)")
        self.btn_view_orig.setObjectName("pill_button")
        self.btn_view_orig.setCheckable(True)
        self.btn_view_orig.setChecked(True)
        self.btn_view_orig.setCursor(Qt.PointingHandCursor)
        
        self.btn_view_gallery = QPushButton("🗂️ 分层预览 (decomposed)")
        self.btn_view_gallery.setObjectName("pill_button")
        self.btn_view_gallery.setCheckable(True)
        self.btn_view_gallery.setCursor(Qt.PointingHandCursor)
        
        view_tabs_layout.addWidget(self.btn_view_orig)
        view_tabs_layout.addWidget(self.btn_view_gallery)
        view_tabs_layout.addStretch()
        
        self.view_group = QButtonGroup(self.parent)
        self.view_group.addButton(self.btn_view_orig)
        self.view_group.addButton(self.btn_view_gallery)
        self.view_group.setExclusive(True)
        
        self.btn_view_orig.clicked.connect(lambda: self.switch_view(0))
        self.btn_view_gallery.clicked.connect(lambda: self.switch_view(1))
        right_layout.addLayout(view_tabs_layout)
        
        # Stacked pages
        self.preview_stack = QStackedWidget()
        
        # Page 0: Original Preview
        self.orig_preview = ImagePreviewWidget()
        self.preview_stack.addWidget(self.orig_preview)
        
        # Page 1: Layers Gallery Scroll Area
        self.scroll_layers = QScrollArea()
        self.scroll_layers.setWidgetResizable(True)
        self.scroll_layers.setStyleSheet("QScrollArea { border: none; background: transparent; }")
        
        self.scroll_content = QWidget()
        self.scroll_content.setObjectName("scroll_page")
        self.grid_layout = QGridLayout(self.scroll_content)
        self.grid_layout.setContentsMargins(10, 10, 10, 10)
        self.grid_layout.setSpacing(15)
        
        self.scroll_layers.setWidget(self.scroll_content)
        self.preview_stack.addWidget(self.scroll_layers)
        
        right_layout.addWidget(self.preview_stack, 1)
        splitter.addWidget(right_panel)
        
        main_layout.addWidget(splitter, 1)
        
    def create_separator(self):
        line = QFrame()
        line.setFrameShape(QFrame.HLine)
        line.setFrameShadow(QFrame.Sunken)
        line.setStyleSheet("background-color: #2e2e32; max-height: 1px;")
        return line
        
    def change_layers_val(self, val):
        self.lbl_layers_val.setText(str(val))
        
    def change_steps_val(self, val):
        self.lbl_steps_val.setText(str(val))
        
    def change_cfg_val(self, val):
        self.lbl_cfg_val.setText(f"{val/10.0:.1f}")
        
    def choose_image(self):
        file_path, _ = QFileDialog.getOpenFileName(
            self.parent, "选择图片", "", "图片文件 (*.png *.jpg *.jpeg *.bmp)"
        )
        if file_path:
            self.selected_img_path = file_path
            self.orig_preview.set_image(file_path)
            
            img = QImage(file_path)
            if not img.isNull():
                self.lbl_image_info.setText(
                    f"图片路径: {os.path.basename(file_path)}\n分辨率: {img.width()} x {img.height()}"
                )
                self.lbl_status.setText("状态: 已载入图片，可配置参数并运行分层。")
                self.btn_run.setEnabled(True)
                
                # Reset actions
                self.btn_export_pptx.setEnabled(False)
                self.btn_export_zip.setEnabled(False)
                self.btn_export_psd.setEnabled(False)
                
                # Switch tab
                self.btn_view_orig.setChecked(True)
                self.switch_view(0)
            else:
                QMessageBox.critical(self.parent, "错误", "图片解析失败，请检查格式是否支持。")
                self.lbl_image_info.setText("加载失败。")
                self.btn_run.setEnabled(False)
                
    def switch_view(self, index):
        self.preview_stack.setCurrentIndex(index)
        
    def run_layering(self):
        if not self.selected_img_path:
            return
            
        layers_count = self.slider_layers.value()
        steps = self.slider_steps.value()
        cfg = self.slider_cfg.value() / 10.0
        prompt = self.txt_prompt.text().strip()
        
        # Disable UI
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)
        self.btn_run.setEnabled(False)
        self.btn_select_image.setEnabled(False)
        
        self.btn_export_pptx.setEnabled(False)
        self.btn_export_zip.setEnabled(False)
        self.btn_export_psd.setEnabled(False)
        
        # Clear existing grid
        while self.grid_layout.count():
            item = self.grid_layout.takeAt(0)
            widget = item.widget()
            if widget:
                widget.deleteLater()
                
        self.lbl_status.setText("状态: 正在建立 AI 分层通道...")
        
        low_vram = self.chk_low_vram.isChecked()
        
        self.worker = QwenLayeredWorker(
            img_path=self.selected_img_path,
            layers_count=layers_count,
            num_steps=steps,
            scale=cfg,
            prompt=prompt,
            seed=-1,
            low_vram=low_vram
        )
        self.worker.progress.connect(self.on_worker_progress)
        self.worker.finished.connect(self.on_worker_finished)
        self.worker.start()
        
    def on_worker_progress(self, msg):
        self.lbl_status.setText(f"状态: {msg}")
        
    def on_worker_finished(self, success, layer_paths, pptx_path, zip_path, psd_path, err_msg):
        self.progress_bar.setVisible(False)
        self.btn_run.setEnabled(True)
        self.btn_select_image.setEnabled(True)
        
        if success:
            self.lbl_status.setText("状态: 智能分层处理完成！")
            self.layer_paths = layer_paths
            self.pptx_path = pptx_path
            self.zip_path = zip_path
            self.psd_path = psd_path
            
            # Enable exports
            self.btn_export_pptx.setEnabled(True)
            self.btn_export_zip.setEnabled(True)
            self.btn_export_psd.setEnabled(True)
            
            # Populate grid
            col_count = 3
            for idx, path in enumerate(self.layer_paths):
                card = LayerCardWidget(idx + 1, path, self.on_layer_preview, self.on_layer_save)
                row = idx // col_count
                col = idx % col_count
                self.grid_layout.addWidget(card, row, col)
                
            # Switch views to decomposed gallery
            self.btn_view_gallery.setChecked(True)
            self.switch_view(1)
            
            QMessageBox.information(self.parent, "分层成功", f"图像已成功解耦为 {len(layer_paths)} 个透明图层，您可以通过下方按钮导出整图工程，或单独导出/预览各图层。")
        else:
            self.lbl_status.setText("状态: 图像分层失败。")
            QMessageBox.critical(self.parent, "分层错误", f"Qwen 模型分层计算出错：\n{err_msg}")
            
    def on_layer_preview(self, img_path, idx):
        dialog = LayerPreviewDialog(img_path, idx, self.parent)
        dialog.exec()
        
    def on_layer_save(self, img_path, idx):
        save_path, _ = QFileDialog.getSaveFileName(
            self.parent, f"保存图层 {idx} 为透明 PNG", f"layer_{idx}.png", "透明图像 (*.png)"
        )
        if save_path:
            try:
                shutil.copy2(img_path, save_path)
                QMessageBox.information(self.parent, "成功", f"图层 {idx} 已成功导出保存！")
            except Exception as e:
                QMessageBox.critical(self.parent, "错误", f"另存为文件失败：\n{e}")
                
    def save_pptx(self):
        if not self.pptx_path or not os.path.exists(self.pptx_path):
            return
        save_path, _ = QFileDialog.getSaveFileName(
            self.parent, "另存为幻灯片 PPTX 文件", "decomposed_image.pptx", "幻灯片文件 (*.pptx)"
        )
        if save_path:
            try:
                shutil.copy2(self.pptx_path, save_path)
                QMessageBox.information(self.parent, "导出成功", f"幻灯片已导出到：\n{save_path}")
            except Exception as e:
                QMessageBox.critical(self.parent, "错误", f"保存失败：\n{e}")
                
    def save_zip(self):
        if not self.zip_path or not os.path.exists(self.zip_path):
            return
        save_path, _ = QFileDialog.getSaveFileName(
            self.parent, "另存为图层 ZIP 压缩包", "decomposed_layers.zip", "压缩包文件 (*.zip)"
        )
        if save_path:
            try:
                shutil.copy2(self.zip_path, save_path)
                QMessageBox.information(self.parent, "导出成功", f"压缩包已保存到：\n{save_path}")
            except Exception as e:
                QMessageBox.critical(self.parent, "错误", f"保存失败：\n{e}")
                
    def save_psd(self):
        if not self.psd_path or not os.path.exists(self.psd_path):
            return
        save_path, _ = QFileDialog.getSaveFileName(
            self.parent, "另存为 Photoshop PSD 工程图层文件", "decomposed_image.psd", "Photoshop 文件 (*.psd)"
        )
        if save_path:
            try:
                shutil.copy2(self.psd_path, save_path)
                QMessageBox.information(self.parent, "导出成功", f"PSD 工程文件已保存到：\n{save_path}")
            except Exception as e:
                QMessageBox.critical(self.parent, "错误", f"保存失败：\n{e}")
